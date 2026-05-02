"""Build the GharSetu capstone defense presentation.

Run: python3 scripts/build_presentation.py
Output: Akshit_Thakur_Capstone_Presentation.pptx (16:9, 22 slides)

Re-runnable. Overwrites the output each run.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree


REPO_ROOT = Path(__file__).resolve().parent.parent
OUTPUT = REPO_ROOT / "Akshit_Thakur_Capstone_Presentation.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

CREAM = RGBColor(0xFF, 0xFB, 0xF5)
INK = RGBColor(0x1A, 0x1A, 0x1A)
INK_SOFT = RGBColor(0x3A, 0x3A, 0x3A)
MUTED = RGBColor(0x55, 0x55, 0x55)
SAFFRON = RGBColor(0xE2, 0x6A, 0x2C)
DEEP = RGBColor(0x8A, 0x3F, 0x18)
GREEN = RGBColor(0x1F, 0x7A, 0x3D)
NAVY = RGBColor(0x1A, 0x35, 0x5E)
PALE = RGBColor(0xF5, 0xEC, 0xDC)
LINE = RGBColor(0xC9, 0xB7, 0x9A)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG = RGBColor(0x1F, 0x24, 0x2C)
CODE_TXT = RGBColor(0xE6, 0xE6, 0xE6)
CODE_KEY = RGBColor(0xF5, 0xC4, 0x6B)
CODE_STR = RGBColor(0x9D, 0xD8, 0xA2)

BODY_FONT = "Calibri"
CODE_FONT = "Consolas"
HINDI_FONT = "Mangal"


def set_solid_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_fill(shape):
    shape.fill.background()


def no_line(shape):
    shape.line.fill.background()


def set_line(shape, color, width_pt=0.75):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE):
    rect = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        no_fill(rect)
    else:
        set_solid_fill(rect, fill)
    if line is None:
        no_line(rect)
    else:
        set_line(rect, line, line_w)
    rect.shadow.inherit = False
    return rect


def add_text(slide, x, y, w, h, text, *, font=BODY_FONT, size=18, bold=False,
             color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_lines(slide, x, y, w, h, lines, *, font=BODY_FONT, size=18, color=INK,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.2,
              space_after=4, bold=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        if isinstance(line, str):
            run = p.add_run()
            run.text = line
            run.font.name = font
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = bold
        else:
            for piece in line:
                run = p.add_run()
                run.text = piece.get("text", "")
                run.font.name = piece.get("font", font)
                run.font.size = Pt(piece.get("size", size))
                run.font.bold = piece.get("bold", False)
                run.font.italic = piece.get("italic", False)
                run.font.color.rgb = piece.get("color", color)
    return tb


def page_background(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    set_solid_fill(bg, CREAM)
    no_line(bg)
    bg.shadow.inherit = False
    return bg


def accent_bar(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.18))
    set_solid_fill(bar, SAFFRON)
    no_line(bar)
    bar.shadow.inherit = False
    return bar


def footer(slide, slide_no, total=22):
    add_text(
        slide, Inches(0.45), Inches(7.10), Inches(8.0), Inches(0.3),
        "Akshit Thakur  ·  GharSetu Capstone  ·  2026",
        size=10, color=MUTED, align=PP_ALIGN.LEFT,
    )
    add_text(
        slide, Inches(11.5), Inches(7.10), Inches(1.4), Inches(0.3),
        f"{slide_no} / {total}",
        size=10, color=MUTED, align=PP_ALIGN.RIGHT,
    )
    line = slide.shapes.add_connector(1, Inches(0.45), Inches(7.05), Inches(12.88), Inches(7.05))
    line.line.color.rgb = LINE
    line.line.width = Pt(0.5)


def slide_title(slide, title, subtitle=None):
    add_text(
        slide, Inches(0.6), Inches(0.45), Inches(12.1), Inches(0.7),
        title, size=32, bold=True, color=INK, anchor=MSO_ANCHOR.TOP,
    )
    if subtitle:
        add_text(
            slide, Inches(0.6), Inches(1.05), Inches(12.1), Inches(0.45),
            subtitle, size=16, color=MUTED, anchor=MSO_ANCHOR.TOP,
        )
    underline = slide.shapes.add_connector(
        1, Inches(0.6),
        Inches(1.55) if subtitle else Inches(1.20),
        Inches(2.4),
        Inches(1.55) if subtitle else Inches(1.20),
    )
    underline.line.color.rgb = SAFFRON
    underline.line.width = Pt(2.5)


def base_slide(prs, slide_no, title=None, subtitle=None):
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)
    page_background(s)
    accent_bar(s)
    if title is not None:
        slide_title(s, title, subtitle)
    footer(s, slide_no)
    return s


def add_arrow(slide, x1, y1, x2, y2, color=INK_SOFT, width_pt=1.5):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width_pt)
    ln = line.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("h", "med")
    return line


def card(slide, x, y, w, h, *, fill=CARD_BG, line_color=LINE, line_w=0.75):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.adjustments[0] = 0.06
    set_solid_fill(box, fill)
    set_line(box, line_color, line_w)
    box.shadow.inherit = False
    return box


def code_block(slide, x, y, w, h, lines, *, size=12):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.adjustments[0] = 0.04
    set_solid_fill(box, CODE_BG)
    no_line(box)
    box.shadow.inherit = False
    pad = Inches(0.18)
    tb = slide.shapes.add_textbox(x + pad, y + pad, w - 2 * pad, h - 2 * pad)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.15
        if isinstance(line, str):
            r = p.add_run()
            r.text = line
            r.font.name = CODE_FONT
            r.font.size = Pt(size)
            r.font.color.rgb = CODE_TXT
        else:
            for piece in line:
                r = p.add_run()
                r.text = piece.get("text", "")
                r.font.name = CODE_FONT
                r.font.size = Pt(piece.get("size", size))
                r.font.color.rgb = piece.get("color", CODE_TXT)
                r.font.bold = piece.get("bold", False)
    return box


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def slide_title_page(prs):
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)
    page_background(s)

    # Saffron block on left
    block = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.55), SLIDE_H)
    set_solid_fill(block, SAFFRON)
    no_line(block)
    block.shadow.inherit = False

    # Wordmark
    add_text(
        s, Inches(1.0), Inches(1.2), Inches(11.5), Inches(1.6),
        "GharSetu", size=96, bold=True, color=INK,
    )
    # Hindi tagline
    add_text(
        s, Inches(1.05), Inches(2.85), Inches(11.5), Inches(0.7),
        "Apna kamra, apne sheher mein.",
        size=28, italic=True, color=DEEP, font=HINDI_FONT,
    )
    # English subtitle
    add_text(
        s, Inches(1.05), Inches(3.55), Inches(11.5), Inches(0.7),
        "A Localized PG and Room Rental Platform for University Students,",
        size=20, color=INK_SOFT,
    )
    add_text(
        s, Inches(1.05), Inches(3.95), Inches(11.5), Inches(0.7),
        "with ONDC Connectivity",
        size=20, color=INK_SOFT,
    )

    # Divider
    div = s.shapes.add_connector(1, Inches(1.05), Inches(4.95), Inches(7.0), Inches(4.95))
    div.line.color.rgb = SAFFRON
    div.line.width = Pt(2.5)

    add_text(
        s, Inches(1.05), Inches(5.10), Inches(11.5), Inches(0.5),
        "Capstone Defense  ·  May 2026",
        size=18, bold=True, color=SAFFRON,
    )
    add_text(
        s, Inches(1.05), Inches(5.55), Inches(11.5), Inches(0.5),
        "Akshit Thakur  ·  B.Tech CSE (Cybersecurity)",
        size=16, color=INK,
    )
    add_text(
        s, Inches(1.05), Inches(5.95), Inches(11.5), Inches(0.45),
        "Yogananda School of AI, Computers and Data Sciences",
        size=14, color=MUTED,
    )
    add_text(
        s, Inches(1.05), Inches(6.30), Inches(11.5), Inches(0.45),
        "Shoolini University, Solan, Himachal Pradesh",
        size=14, color=MUTED,
    )

    # Bottom corner mark
    add_text(
        s, Inches(1.05), Inches(6.85), Inches(11.5), Inches(0.4),
        "Apna kamra. Apna sheher. Apna ONDC.",
        size=12, italic=True, color=MUTED,
    )


def slide_problem(prs):
    s = base_slide(prs, 2, "The Problem", "Why a localized rental platform, why now")

    pts = [
        ("Scattered supply.",
         "Listings live across WhatsApp groups, paper notices on hostel walls, and brokers. No single search."),
        ("No honest reviews.",
         "Today's reviews are written by people who never paid rent there. Renters who actually lived in the room have no place to speak."),
        ("Deposit and quality fraud.",
         "Promised amenities vanish on move-in day. Deposits get withheld. Day-scholar and out-of-town students absorb the loss."),
    ]
    y = Inches(1.85)
    for title, body in pts:
        card(s, Inches(0.6), y, Inches(8.7), Inches(1.35), fill=CARD_BG, line_color=LINE)
        add_text(s, Inches(0.85), y + Inches(0.18), Inches(8.2), Inches(0.45),
                 title, size=20, bold=True, color=DEEP)
        add_text(s, Inches(0.85), y + Inches(0.62), Inches(8.2), Inches(0.7),
                 body, size=14, color=INK_SOFT)
        y += Inches(1.55)

    # Quote panel from idea.md
    card(s, Inches(9.55), Inches(1.85), Inches(3.3), Inches(4.55), fill=PALE, line_color=SAFFRON, line_w=1.0)
    add_text(s, Inches(9.75), Inches(2.0), Inches(3.0), Inches(0.4),
             "From the founder",
             size=11, bold=True, color=DEEP)
    add_text(s, Inches(9.75), Inches(2.4), Inches(3.0), Inches(3.8),
             "“Ye chiz bohot pareshan kerte hh "
             "day scholar students ko, islie maine "
             "ye topic socha.”",
             size=14, italic=True, color=INK, font=HINDI_FONT)
    add_text(s, Inches(9.75), Inches(5.85), Inches(3.0), Inches(0.5),
             "— Akshit Thakur, idea.md",
             size=11, color=MUTED)


def slide_insight(prs):
    s = base_slide(prs, 3, "The Insight", "Three words that scope every product decision")

    cols = [
        ("Localized", "One radius. One campus. One city. The defaults are tuned for the place the student actually lives.", SAFFRON),
        ("Federated", "ONDC means one listing reaches every buyer-app on the network. The owner lists once, gets discovered everywhere.", NAVY),
        ("Verified", "Reviews carry a renter badge if the writer actually paid rent or the owner marked them. Outsiders can still speak; readers can still tell.", GREEN),
    ]
    x = Inches(0.6)
    w = Inches(4.04)
    for title, body, color in cols:
        card(s, x, Inches(2.3), w, Inches(3.7), fill=CARD_BG, line_color=LINE)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.3), w, Inches(0.12))
        set_solid_fill(bar, color)
        no_line(bar)
        bar.shadow.inherit = False
        add_text(s, x + Inches(0.3), Inches(2.7), w - Inches(0.6), Inches(0.7),
                 title, size=28, bold=True, color=color)
        add_text(s, x + Inches(0.3), Inches(3.5), w - Inches(0.6), Inches(2.0),
                 body, size=15, color=INK_SOFT, align=PP_ALIGN.LEFT)
        x += w + Inches(0.16)

    add_text(s, Inches(0.6), Inches(6.20), Inches(12.1), Inches(0.6),
             "Every feature ships only if it serves Localized + Federated + Verified.",
             size=15, italic=True, color=MUTED, align=PP_ALIGN.CENTER)


def slide_vision(prs):
    s = base_slide(prs, 4, "Vision", "What a student should feel on day one")

    big_lines = [
        "Search a room near your campus in 30 seconds.",
        "Trust because the reviews are from people who actually paid rent.",
        "Pay safely on the platform.",
        "List once, reach every ONDC buyer-app in the country.",
    ]
    y = Inches(2.0)
    colors = [SAFFRON, GREEN, NAVY, DEEP]
    for line, color in zip(big_lines, colors):
        # bullet circle
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85), y + Inches(0.12), Inches(0.22), Inches(0.22))
        set_solid_fill(dot, color)
        no_line(dot)
        dot.shadow.inherit = False
        add_text(s, Inches(1.35), y, Inches(11.4), Inches(0.7),
                 line, size=24, bold=False, color=INK)
        y += Inches(0.95)

    # Closing line
    add_text(s, Inches(0.6), Inches(6.20), Inches(12.1), Inches(0.6),
             "Built for a slow phone, a flaky 2G connection, and a small-town budget.",
             size=15, italic=True, color=MUTED, align=PP_ALIGN.CENTER)


def slide_users(prs):
    s = base_slide(prs, 5, "Users and Roles", "Three personas, one shared trust loop")

    personas = [
        {
            "title": "Student",
            "color": SAFFRON,
            "icon": "S",
            "actions": [
                "Search by city, budget, gender pref.",
                "Request a visit or reserve a room.",
                "Pay rent on the platform.",
                "Leave feedback (verified badge).",
                "Verify identity via DigiLocker.",
            ],
        },
        {
            "title": "Owner",
            "color": NAVY,
            "icon": "O",
            "actions": [
                "Create / edit / pause listings.",
                "Receive booking inbox.",
                "Accept or decline visit requests.",
                "Mark a tenant as renter manually.",
                "Track captured payments.",
            ],
        },
        {
            "title": "Admin",
            "color": GREEN,
            "icon": "A",
            "actions": [
                "Live SIEM dashboard (audit log).",
                "Filter by actor / action / entity.",
                "Spot abuse, rate-limit hits, errors.",
                "Inspect ONDC message timeline.",
                "Read-only by design.",
            ],
        },
    ]
    x = Inches(0.6)
    w = Inches(4.04)
    for p in personas:
        card(s, x, Inches(1.85), w, Inches(4.65), fill=CARD_BG, line_color=LINE)
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), Inches(2.05), Inches(0.7), Inches(0.7))
        set_solid_fill(circle, p["color"])
        no_line(circle)
        circle.shadow.inherit = False
        ic = circle.text_frame
        ic.margin_left = Emu(0); ic.margin_right = Emu(0)
        ic.margin_top = Emu(0); ic.margin_bottom = Emu(0)
        ic.vertical_anchor = MSO_ANCHOR.MIDDLE
        ic.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = ic.paragraphs[0].add_run()
        run.text = p["icon"]
        run.font.name = BODY_FONT
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = CARD_BG
        add_text(s, x + Inches(1.1), Inches(2.15), w - Inches(1.3), Inches(0.55),
                 p["title"], size=24, bold=True, color=p["color"])

        ay = Inches(2.95)
        for a in p["actions"]:
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), ay + Inches(0.13), Inches(0.10), Inches(0.10))
            set_solid_fill(dot, p["color"])
            no_line(dot)
            dot.shadow.inherit = False
            add_text(s, x + Inches(0.55), ay, w - Inches(0.7), Inches(0.5),
                     a, size=13.5, color=INK_SOFT)
            ay += Inches(0.62)
        x += w + Inches(0.16)


def slide_built(prs):
    s = base_slide(prs, 6, "What's Been Built (MLP)", "End-to-end, deployable, tested")

    # Stat cards across top
    stats = [
        ("36", "HTTP routes"),
        ("10", "DB tables"),
        ("9", "ONDC actions"),
        ("2", "languages"),
        ("10/10", "smoke tests pass"),
    ]
    x = Inches(0.6)
    w = Inches(2.4)
    for value, label in stats:
        card(s, x, Inches(1.75), w, Inches(1.35), fill=PALE, line_color=SAFFRON, line_w=1.0)
        add_text(s, x, Inches(1.85), w, Inches(0.7),
                 value, size=32, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(2.55), w, Inches(0.5),
                 label, size=12, color=MUTED, align=PP_ALIGN.CENTER)
        x += w + Inches(0.08)

    # Feature columns
    add_text(s, Inches(0.6), Inches(3.35), Inches(6.0), Inches(0.5),
             "Surface", size=18, bold=True, color=NAVY)
    feats_left = [
        "Server-rendered EJS pages, mobile first.",
        "Search with city, budget, gender, amenities.",
        "Listings CRUD with image upload (sharp WebP).",
        "Bookings: visit + reserve, owner inbox.",
        "Razorpay rent payments (real-shape sim).",
        "DigiLocker KYC (real-shape sim).",
        "Verified-renter feedback engine.",
    ]
    add_lines(s, Inches(0.6), Inches(3.85), Inches(6.0), Inches(3.2),
              [f"•  {f}" for f in feats_left], size=14, color=INK_SOFT)

    add_text(s, Inches(6.9), Inches(3.35), Inches(6.0), Inches(0.5),
             "Platform", size=18, bold=True, color=NAVY)
    feats_right = [
        "ONDC Beckn 1.2 endpoints (search→support).",
        "Bilingual EN + HI, RTL-ready, locale picker.",
        "WCAG 2.2 AAA: skip-link, focus, ARIA, contrast.",
        "Service worker shell + last 20 listings cache.",
        "Pino structured logs to stdout.",
        "Admin SIEM dashboard with live SSE feed.",
        "Cloud Run deploy via one shell script.",
    ]
    add_lines(s, Inches(6.9), Inches(3.85), Inches(6.0), Inches(3.2),
              [f"•  {f}" for f in feats_right], size=14, color=INK_SOFT)


def slide_demo(prs):
    s = base_slide(prs, 7, "Live Demo", "What you'd see on screen")

    # Five panels in a strip
    panels = [
        ("Home /", "Hero, search bar, featured listings near Solan. 12KB HTML, 25KB CSS."),
        ("Search /search", "Filters left, list + map toggle right. JSON API at /api/search."),
        ("Listing /listings/:id", "Photos, map, amenities, verified-renter feedback first."),
        ("Book /bookings", "Visit or reserve. Owner gets it in their dashboard inbox."),
        ("Pay /pay/:id", "Razorpay-style order created. Webhook flips status to captured."),
        ("Admin /admin", "Live audit feed, filter by action / actor, ONDC traffic visible."),
    ]
    cols = 3
    pw = Inches(4.04)
    ph = Inches(2.0)
    gap_x = Inches(0.16)
    gap_y = Inches(0.20)
    sx = Inches(0.6)
    sy = Inches(1.85)
    for i, (h, body) in enumerate(panels):
        r = i // cols
        c = i % cols
        x = sx + c * (pw + gap_x)
        y = sy + r * (ph + gap_y)
        card(s, x, y, pw, ph, fill=CARD_BG, line_color=LINE)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, pw, Inches(0.06))
        set_solid_fill(bar, SAFFRON)
        no_line(bar)
        bar.shadow.inherit = False
        add_text(s, x + Inches(0.25), y + Inches(0.20), pw - Inches(0.5), Inches(0.45),
                 h, size=16, bold=True, color=INK)
        add_text(s, x + Inches(0.25), y + Inches(0.70), pw - Inches(0.5), Inches(1.20),
                 body, size=12, color=INK_SOFT)

    # Demo creds bar
    card(s, Inches(0.6), Inches(6.30), Inches(12.28), Inches(0.55), fill=PALE, line_color=SAFFRON, line_w=1.0)
    add_text(
        s, Inches(0.8), Inches(6.36), Inches(12.0), Inches(0.45),
        "Seeded credentials   "
        "student@gharsetu.local / Student@2026!     "
        "owner@gharsetu.local / Owner@2026!     "
        "admin@gharsetu.local / Admin@2026!",
        size=12, color=INK,
    )


def slide_architecture(prs):
    s = base_slide(prs, 8, "System Architecture", "One container, one process, zero idle cost")

    # Browser
    browser = card(s, Inches(0.6), Inches(2.6), Inches(2.4), Inches(1.4), fill=CARD_BG, line_color=NAVY, line_w=1.0)
    add_text(s, Inches(0.6), Inches(2.75), Inches(2.4), Inches(0.5),
             "Browser", size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(3.15), Inches(2.4), Inches(0.4),
             "(mobile, 2G ok)", size=12, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(3.45), Inches(2.4), Inches(0.4),
             "EJS + 0 framework JS", size=11, color=INK_SOFT, align=PP_ALIGN.CENTER)

    # Cloud Run server box
    cr = card(s, Inches(4.4), Inches(2.0), Inches(4.5), Inches(2.7), fill=PALE, line_color=SAFFRON, line_w=1.5)
    add_text(s, Inches(4.4), Inches(2.10), Inches(4.5), Inches(0.5),
             "Cloud Run (asia-south1)", size=11, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
    add_text(s, Inches(4.4), Inches(2.45), Inches(4.5), Inches(0.5),
             "Fastify 5  +  Node 22 LTS", size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)

    # Inside CR: stack
    inner_x = Inches(4.7)
    inner_w = Inches(3.9)
    for i, label in enumerate(["EJS SSR", "zod validation", "JWT cookie auth", "pino JSON logs"]):
        y = Inches(3.05 + i * 0.36)
        b = add_rect(s, inner_x, y, inner_w, Inches(0.32), fill=CARD_BG, line=LINE)
        add_text(s, inner_x, y + Inches(0.04), inner_w, Inches(0.30),
                 label, size=12, color=INK_SOFT, align=PP_ALIGN.CENTER)

    # Storage on right
    db_box = card(s, Inches(10.0), Inches(2.0), Inches(2.9), Inches(1.25), fill=CARD_BG, line_color=GREEN, line_w=1.0)
    add_text(s, Inches(10.0), Inches(2.10), Inches(2.9), Inches(0.5),
             "SQLite", size=18, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(s, Inches(10.0), Inches(2.55), Inches(2.9), Inches(0.4),
             "/tmp/gharsetu.db", size=12, color=MUTED, align=PP_ALIGN.CENTER, font=CODE_FONT)
    add_text(s, Inches(10.0), Inches(2.85), Inches(2.9), Inches(0.4),
             "ephemeral, seeded on cold start", size=11, color=INK_SOFT, align=PP_ALIGN.CENTER)

    up_box = card(s, Inches(10.0), Inches(3.45), Inches(2.9), Inches(1.25), fill=CARD_BG, line_color=GREEN, line_w=1.0)
    add_text(s, Inches(10.0), Inches(3.55), Inches(2.9), Inches(0.5),
             "/tmp/uploads", size=16, bold=True, color=GREEN, align=PP_ALIGN.CENTER, font=CODE_FONT)
    add_text(s, Inches(10.0), Inches(4.00), Inches(2.9), Inches(0.6),
             "sharp → WebP @ q=78", size=12, color=INK_SOFT, align=PP_ALIGN.CENTER)

    # Arrows
    add_arrow(s, Inches(3.0), Inches(3.3), Inches(4.4), Inches(3.3), color=INK_SOFT)
    add_arrow(s, Inches(8.9), Inches(2.6), Inches(10.0), Inches(2.6), color=INK_SOFT)
    add_arrow(s, Inches(8.9), Inches(4.05), Inches(10.0), Inches(4.05), color=INK_SOFT)

    # Side peers (bottom row)
    peers = [
        ("ONDC peers (BAP / BPP)", "POST /ondc/v1/search...support  →  ACK + async on_*", NAVY),
        ("Razorpay (sim)", "Order create  +  webhook HMAC-SHA256  →  payment.captured", DEEP),
        ("DigiLocker (sim)", "OAuth 2.0 PKCE flow  →  callback  →  kyc_verified=1", GREEN),
    ]
    x = Inches(0.6)
    w = Inches(4.10)
    for title, body, color in peers:
        card(s, x, Inches(5.20), w, Inches(1.55), fill=CARD_BG, line_color=color, line_w=1.0)
        add_text(s, x + Inches(0.2), Inches(5.30), w - Inches(0.4), Inches(0.4),
                 title, size=13, bold=True, color=color)
        add_text(s, x + Inches(0.2), Inches(5.65), w - Inches(0.4), Inches(1.0),
                 body, size=11, color=INK_SOFT, font=CODE_FONT)
        x += w + Inches(0.04)


def slide_data_model(prs):
    s = base_slide(prs, 9, "Data Model", "Ten tables. The trust badge is the centrepiece")

    tables = [
        ("users", "id, email, role, kyc_verified, kyc_method", NAVY),
        ("listings", "id, owner_id, city, lat, lng, rent_monthly, status", NAVY),
        ("listing_images", "id, listing_id, url, position", NAVY),
        ("bookings", "id, listing_id, student_id, type, status", NAVY),
        ("payments", "id, listing_id, payer_id, rzp_order_id, status", DEEP),
        ("renter_records", "id, listing_id, student_id, source, active", SAFFRON),
        ("feedback", "id, listing_id, author_id, rating, is_verified_renter", GREEN),
        ("audit_log", "id, actor_id, action, entity, payload", MUTED),
        ("ondc_messages", "id, txn_id, message_id, action, direction", NAVY),
        ("sessions", "jti, user_id, revoked, expires_at", MUTED),
    ]
    cols = 5
    tw = Inches(2.43)
    th = Inches(1.20)
    gx = Inches(0.10)
    gy = Inches(0.20)
    sx = Inches(0.6)
    sy = Inches(1.85)
    for i, (name, body, color) in enumerate(tables):
        r = i // cols
        c = i % cols
        x = sx + c * (tw + gx)
        y = sy + r * (th + gy)
        card(s, x, y, tw, th, fill=CARD_BG, line_color=color, line_w=1.0)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, tw, Inches(0.05))
        set_solid_fill(bar, color); no_line(bar); bar.shadow.inherit = False
        add_text(s, x + Inches(0.12), y + Inches(0.10), tw - Inches(0.24), Inches(0.4),
                 name, size=14, bold=True, color=color, font=CODE_FONT)
        add_text(s, x + Inches(0.12), y + Inches(0.45), tw - Inches(0.24), Inches(0.75),
                 body, size=10, color=INK_SOFT, font=CODE_FONT)

    # Flow band
    band_y = Inches(4.55)
    add_text(s, Inches(0.6), band_y, Inches(12.1), Inches(0.45),
             "Trust flow", size=16, bold=True, color=DEEP)

    chain = [
        ("payments", DEEP),
        ("renter_records", SAFFRON),
        ("feedback", GREEN),
    ]
    bx = Inches(0.8)
    bw = Inches(2.7)
    bh = Inches(0.65)
    by = Inches(5.05)
    boxes = []
    for name, color in chain:
        r = card(s, bx, by, bw, bh, fill=PALE, line_color=color, line_w=1.0)
        add_text(s, bx, by + Inches(0.16), bw, Inches(0.4),
                 name, size=14, bold=True, color=color, font=CODE_FONT, align=PP_ALIGN.CENTER)
        boxes.append((bx, by, bw, bh))
        bx += bw + Inches(0.55)

    # Arrows between boxes
    for i in range(len(boxes) - 1):
        x1 = boxes[i][0] + boxes[i][2]
        x2 = boxes[i + 1][0]
        ym = boxes[i][1] + Inches(0.32)
        add_arrow(s, x1, ym, x2, ym, color=INK_SOFT, width_pt=1.5)

    # Owner-marked side input
    om = card(s, Inches(9.6), Inches(4.85), Inches(3.3), Inches(1.05), fill=CARD_BG, line_color=SAFFRON, line_w=1.0)
    add_text(s, Inches(9.6), Inches(4.92), Inches(3.3), Inches(0.4),
             "Owner UI", size=12, bold=True, color=SAFFRON, align=PP_ALIGN.CENTER)
    add_text(s, Inches(9.6), Inches(5.28), Inches(3.3), Inches(0.6),
             "POST /owner/listings/:lid/renters\nsource = owner_marked",
             size=10, color=INK_SOFT, font=CODE_FONT, align=PP_ALIGN.CENTER)
    add_arrow(s, Inches(9.6), Inches(5.37), Inches(8.55), Inches(5.37),
              color=SAFFRON, width_pt=1.5)

    # Caption
    add_text(s, Inches(0.6), Inches(6.10), Inches(12.1), Inches(0.5),
             "renter_records is the single source of truth for the verified-renter badge.",
             size=12, italic=True, color=MUTED, align=PP_ALIGN.CENTER)


def slide_trust(prs):
    s = base_slide(prs, 10, "The Verified-Renter Trust Mechanism",
                   "Outsiders may speak. Readers can tell who actually lived there.")

    # Logical OR
    code_block(
        s, Inches(0.6), Inches(1.90), Inches(12.28), Inches(1.95),
        [
            [
                {"text": "is_verified_renter", "color": CODE_KEY, "bold": True},
                {"text": "  =  EXISTS ("},
            ],
            [
                {"text": "  SELECT 1 FROM renter_records r"},
            ],
            [
                {"text": "   WHERE r.listing_id = :listing_id"},
            ],
            [
                {"text": "     AND r.student_id = :author_id"},
            ],
            [
                {"text": "     AND r.active     = 1"},
            ],
            [
                {"text": "     AND r.source     IN ("},
                {"text": "'platform_payment'", "color": CODE_STR},
                {"text": ", "},
                {"text": "'owner_marked'", "color": CODE_STR},
                {"text": ")"},
            ],
            [{"text": ")"}],
        ],
        size=14,
    )

    # Two paths
    paths = [
        ("Path A  ·  platform_payment",
         "Razorpay webhook HMAC-verifies, marks payment captured, "
         "auto-inserts a renter_record. The student gets the badge automatically.",
         DEEP),
        ("Path B  ·  owner_marked",
         "Owner taps Mark Renter for a tenant who paid offline. "
         "Same renter_record is written. Same badge surfaces on review.",
         SAFFRON),
    ]
    x = Inches(0.6)
    w = Inches(6.05)
    for title, body, color in paths:
        card(s, x, Inches(4.10), w, Inches(1.65), fill=CARD_BG, line_color=color, line_w=1.0)
        add_text(s, x + Inches(0.25), Inches(4.20), w - Inches(0.5), Inches(0.45),
                 title, size=15, bold=True, color=color)
        add_text(s, x + Inches(0.25), Inches(4.65), w - Inches(0.5), Inches(1.05),
                 body, size=12, color=INK_SOFT)
        x += w + Inches(0.18)

    # Why it matters
    card(s, Inches(0.6), Inches(5.90), Inches(12.28), Inches(0.85),
         fill=PALE, line_color=GREEN, line_w=1.0)
    add_text(s, Inches(0.85), Inches(6.00), Inches(11.7), Inches(0.4),
             "Why this matters", size=13, bold=True, color=GREEN)
    add_text(s, Inches(0.85), Inches(6.32), Inches(11.7), Inches(0.5),
             "Outsider reviews still count for visibility, but the reader sees a coloured badge "
             "next to verified ones. No moderation team needed.",
             size=12, color=INK_SOFT)


def slide_ondc(prs):
    s = base_slide(prs, 11, "ONDC Beckn Simulation", "Nine inbound actions, async on_* callbacks, full message log")

    # Sequence diagram
    actors = [
        ("Buyer App", NAVY, Inches(1.4)),
        ("GharSetu /ondc/v1/search", SAFFRON, Inches(5.2)),
        ("ondc_messages", GREEN, Inches(9.8)),
    ]
    for label, color, x in actors:
        card(s, x, Inches(1.85), Inches(2.6), Inches(0.55), fill=color, line_color=color)
        tb = add_text(s, x, Inches(1.93), Inches(2.6), Inches(0.4),
                      label, size=12, bold=True, color=CARD_BG, align=PP_ALIGN.CENTER, font=CODE_FONT)
        # vertical lifeline
        line = s.shapes.add_connector(1, x + Inches(1.3), Inches(2.40), x + Inches(1.3), Inches(4.55))
        line.line.color.rgb = LINE
        line.line.width = Pt(0.75)

    # Sequence steps
    steps = [
        (Inches(2.7), Inches(2.65), 0, 1, "search { context, message.intent }"),
        (Inches(2.7), Inches(3.05), 1, 0, "200 ACK"),
        (Inches(2.7), Inches(3.45), 1, 2, "INSERT direction='in'"),
        (Inches(2.7), Inches(3.85), 1, 2, "INSERT direction='out' (on_search)"),
        (Inches(2.7), Inches(4.25), 1, 0, "POST {bap_uri}/on_search { catalog }"),
    ]
    for tx, ty, src, dst, label in steps:
        x1 = actors[src][2] + Inches(1.3)
        x2 = actors[dst][2] + Inches(1.3)
        add_arrow(s, x1, ty, x2, ty, color=INK_SOFT, width_pt=1.0)
        # label centered
        lx = min(x1, x2) + Inches(0.05)
        lw = abs(x2 - x1) - Inches(0.10) if abs(x2 - x1) > Inches(0.5) else Inches(3.0)
        add_text(s, lx, ty - Inches(0.25), lw, Inches(0.25),
                 label, size=10, color=INK_SOFT, font=CODE_FONT, align=PP_ALIGN.CENTER)

    # Action grid
    add_text(s, Inches(0.6), Inches(4.90), Inches(12.1), Inches(0.4),
             "Implemented inbound endpoints (each replies on_*)",
             size=13, bold=True, color=NAVY)
    actions = ["search", "select", "init", "confirm", "status", "cancel", "update", "rating", "support"]
    cols = 9
    bx = Inches(0.6)
    bw = Inches(1.30)
    bh = Inches(0.55)
    by = Inches(5.30)
    for i, a in enumerate(actions):
        x = bx + i * (bw + Inches(0.06))
        card(s, x, by, bw, bh, fill=PALE, line_color=SAFFRON, line_w=0.75)
        add_text(s, x, by + Inches(0.13), bw, Inches(0.3),
                 a, size=11, bold=True, color=DEEP, align=PP_ALIGN.CENTER, font=CODE_FONT)

    # Registry endpoint + prod note
    add_text(s, Inches(0.6), Inches(6.05), Inches(12.1), Inches(0.4),
             "GET /ondc/v1/lookup  —  subscriber metadata for the registry",
             size=12, color=INK_SOFT, font=CODE_FONT)
    add_text(s, Inches(0.6), Inches(6.40), Inches(12.1), Inches(0.5),
             "Production swap: subscribe with the ONDC registry, sign every request "
             "with Ed25519 + SHA-512, route via the Gateway. Schemas already match.",
             size=12, italic=True, color=MUTED)


def slide_payments(prs):
    s = base_slide(prs, 12, "Razorpay Payment Simulation",
                   "Real-shape Orders API + HMAC-verified webhook")

    # Order shape (left code box)
    code_block(
        s, Inches(0.6), Inches(1.90), Inches(6.05), Inches(2.85),
        [
            [{"text": "// POST /pay/order  -->  Razorpay Orders API shape"}],
            [{"text": "{"}],
            [
                {"text": "  id"},
                {"text": ": "},
                {"text": "\"order_01HXY...\"", "color": CODE_STR},
                {"text": ","},
            ],
            [
                {"text": "  entity"},
                {"text": ": "},
                {"text": "\"order\"", "color": CODE_STR},
                {"text": ","},
            ],
            [
                {"text": "  amount"},
                {"text": ": "},
                {"text": "500000", "color": CODE_KEY},
                {"text": ",  "},
                {"text": "// paise"},
            ],
            [
                {"text": "  currency"},
                {"text": ": "},
                {"text": "\"INR\"", "color": CODE_STR},
                {"text": ","},
            ],
            [
                {"text": "  status"},
                {"text": ": "},
                {"text": "\"created\"", "color": CODE_STR},
                {"text": ","},
            ],
            [
                {"text": "  receipt"},
                {"text": ": "},
                {"text": "\"bk_<bookingId>\"", "color": CODE_STR},
                {"text": ","},
            ],
            [{"text": "  notes: { listing_id, payer_id, payee_id }"}],
            [{"text": "}"}],
        ],
        size=12,
    )

    # Webhook verify
    code_block(
        s, Inches(6.85), Inches(1.90), Inches(6.05), Inches(2.85),
        [
            [{"text": "// POST /pay/webhook  -->  HMAC-SHA256 verify"}],
            [
                {"text": "const "},
                {"text": "expected", "color": CODE_KEY},
                {"text": " = createHmac("},
                {"text": "\"sha256\"", "color": CODE_STR},
                {"text": ","},
            ],
            [{"text": "    config.RZP.WEBHOOK_SECRET)"}],
            [{"text": "  .update(rawBody)"}],
            [
                {"text": "  .digest("},
                {"text": "\"hex\"", "color": CODE_STR},
                {"text": ");"},
            ],
            [
                {"text": "if (!timingSafeEqual(", "color": CODE_TXT},
            ],
            [{"text": "       Buffer.from(expected),"}],
            [{"text": "       Buffer.from(req.headers[\"x-razorpay-signature\"])))"}],
            [
                {"text": "  return reply.code("},
                {"text": "400", "color": CODE_KEY},
                {"text": ").send();"},
            ],
            [{"text": "// captured  -->  insert renter_record (source=platform_payment)"}],
        ],
        size=12,
    )

    # Outcome flow
    flow = [
        ("payment.captured", DEEP),
        ("payments.status = captured", SAFFRON),
        ("renter_records inserted", GREEN),
        ("verified badge unlocked", NAVY),
    ]
    bx = Inches(0.6)
    bw = Inches(2.85)
    bh = Inches(0.65)
    by = Inches(5.20)
    boxes = []
    for label, color in flow:
        card(s, bx, by, bw, bh, fill=CARD_BG, line_color=color, line_w=1.0)
        add_text(s, bx, by + Inches(0.16), bw, Inches(0.4),
                 label, size=12, bold=True, color=color, align=PP_ALIGN.CENTER, font=CODE_FONT)
        boxes.append((bx, by, bw, bh))
        bx += bw + Inches(0.18)

    for i in range(len(boxes) - 1):
        x1 = boxes[i][0] + boxes[i][2]
        x2 = boxes[i + 1][0]
        ym = boxes[i][1] + Inches(0.32)
        add_arrow(s, x1, ym, x2, ym, color=INK_SOFT, width_pt=1.5)

    add_text(s, Inches(0.6), Inches(6.10), Inches(12.28), Inches(0.6),
             "Production swap: real Razorpay keys + real x-razorpay-signature header. "
             "Code path is identical.",
             size=12, italic=True, color=MUTED, align=PP_ALIGN.CENTER)


def slide_kyc(prs):
    s = base_slide(prs, 13, "DigiLocker KYC Simulation",
                   "OAuth 2.0 PKCE flow, sanitized payload only")

    # Sequence between three actors
    actors = [
        ("Student", NAVY, Inches(1.0)),
        ("GharSetu /verify", SAFFRON, Inches(5.2)),
        ("DigiLocker (sim)", GREEN, Inches(9.8)),
    ]
    for label, color, x in actors:
        card(s, x, Inches(1.85), Inches(2.6), Inches(0.55), fill=color, line_color=color)
        add_text(s, x, Inches(1.93), Inches(2.6), Inches(0.4),
                 label, size=12, bold=True, color=CARD_BG, align=PP_ALIGN.CENTER)
        line = s.shapes.add_connector(1, x + Inches(1.3), Inches(2.40), x + Inches(1.3), Inches(5.10))
        line.line.color.rgb = LINE
        line.line.width = Pt(0.75)

    steps = [
        (Inches(2.65), 0, 1, "POST /verify/digilocker/init"),
        (Inches(3.00), 1, 2, "302 redirect (state, code_challenge S256)"),
        (Inches(3.35), 2, 0, "consent screen + confirm"),
        (Inches(3.70), 2, 1, "GET /verify/digilocker/callback?code=SIM_..."),
        (Inches(4.05), 1, 1, "verify state, exchange code, fetch eAadhaar (sim)"),
        (Inches(4.40), 1, 1, "store sanitized payload  →  kyc_verified=1"),
        (Inches(4.75), 1, 0, "redirect to /student/dashboard  +  flash 'verified'"),
    ]
    for ty, src, dst, label in steps:
        if src == dst:
            x1 = actors[src][2] + Inches(1.3)
            self_arc = s.shapes.add_shape(MSO_SHAPE.RIGHT_BRACE, x1, ty - Inches(0.05), Inches(0.30), Inches(0.30))
            no_fill(self_arc); set_line(self_arc, INK_SOFT, 1.0); self_arc.shadow.inherit = False
            add_text(s, x1 + Inches(0.4), ty - Inches(0.10), Inches(5.0), Inches(0.30),
                     label, size=10, color=INK_SOFT, font=CODE_FONT)
        else:
            x1 = actors[src][2] + Inches(1.3)
            x2 = actors[dst][2] + Inches(1.3)
            add_arrow(s, x1, ty, x2, ty, color=INK_SOFT, width_pt=1.0)
            lx = min(x1, x2) + Inches(0.10)
            lw = abs(x2 - x1) - Inches(0.20)
            add_text(s, lx, ty - Inches(0.25), lw, Inches(0.25),
                     label, size=10, color=INK_SOFT, font=CODE_FONT, align=PP_ALIGN.CENTER)

    # Bottom panel
    card(s, Inches(0.6), Inches(5.45), Inches(12.28), Inches(1.30), fill=PALE, line_color=GREEN, line_w=1.0)
    add_text(s, Inches(0.85), Inches(5.55), Inches(11.7), Inches(0.4),
             "Sanitized payload stored", size=13, bold=True, color=GREEN)
    add_text(s, Inches(0.85), Inches(5.85), Inches(11.7), Inches(0.5),
             "{ method: 'digilocker_simulated', verified_at, aadhaar_last4: '1234', name_on_id, dob_year }",
             size=12, color=INK_SOFT, font=CODE_FONT)
    add_text(s, Inches(0.85), Inches(6.20), Inches(11.7), Inches(0.5),
             "Production swap: real CLIENT_ID / CLIENT_SECRET + PKCE verifier persistence. "
             "kyc_verified drives the trust badge across the app.",
             size=11, italic=True, color=MUTED)


def slide_perf(prs):
    s = base_slide(prs, 14, "Performance for Slow Phones",
                   "Designed against a 2G connection on a budget Android")

    # Stat strip
    stats = [
        ("12 KB", "Home HTML", "gzipped"),
        ("25 KB", "Critical CSS", "single file"),
        ("0 KB", "JS framework", "vanilla only"),
        ("78", "WebP quality", "via sharp"),
        ("~5 s", "Cloud Run cold", "min=0 always"),
    ]
    x = Inches(0.6)
    w = Inches(2.4)
    for value, label, sub in stats:
        card(s, x, Inches(1.85), w, Inches(1.55), fill=PALE, line_color=SAFFRON, line_w=1.0)
        add_text(s, x, Inches(1.95), w, Inches(0.7),
                 value, size=28, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(2.65), w, Inches(0.4),
                 label, size=12, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(2.97), w, Inches(0.4),
                 sub, size=10, color=MUTED, align=PP_ALIGN.CENTER)
        x += w + Inches(0.08)

    # Decisions
    add_text(s, Inches(0.6), Inches(3.7), Inches(12.1), Inches(0.5),
             "Engineering decisions that earned those numbers",
             size=16, bold=True, color=NAVY)

    pts = [
        "Server-rendered EJS  —  the page is interactive before any JS parses.",
        "Leaflet + OSM tiles loaded only on map pages, lazy from CDN with SRI.",
        "Sharp re-encodes uploads to WebP at quality 78. Listing page weighs in well under 200 KB.",
        "Service worker caches the app shell + last 20 listings. Repeat search works offline.",
        "Cloud Run gzips by default. min-instances=0  =  zero rupees idle.",
        "No client bundler, no dependency tree to ship. The browser pays only for what it uses.",
    ]
    add_lines(
        s, Inches(0.6), Inches(4.20), Inches(12.28), Inches(2.7),
        [f"•  {p}" for p in pts], size=14, color=INK_SOFT, line_spacing=1.25, space_after=4,
    )


def slide_a11y(prs):
    s = base_slide(prs, 15, "Accessibility (WCAG 2.2 AAA)",
                   "Equal usability for blind, deaf, motor, low-vision, cognitive, ADHD users")

    cols = [
        ("Vision",
         ["Contrast ≥ 7:1 every text element.",
          "Visible focus ring ≥ 3:1 via :focus-visible.",
          "Skip-to-content as first focusable.",
          "ARIA labels on icon-only buttons.",
          "<html lang> set per request."],
         NAVY),
        ("Motion / cognition",
         ["prefers-reduced-motion respected.",
          "forced-colors mode supported.",
          "No autoplay, no time-pressured forms.",
          "Form errors via aria-describedby.",
          "Plain language in EN + HI."],
         SAFFRON),
        ("Locale and direction",
         ["Bilingual EN + HI from launch.",
          "Locale picker in nav, cookie + query param.",
          "RTL-ready: dir attribute per locale.",
          "Mangal / Nirmala UI fallback for Devanagari.",
          "Tested on NVDA / Firefox + VoiceOver / Safari."],
         GREEN),
    ]
    x = Inches(0.6)
    w = Inches(4.04)
    for title, items, color in cols:
        card(s, x, Inches(1.85), w, Inches(4.65), fill=CARD_BG, line_color=LINE)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.85), w, Inches(0.10))
        set_solid_fill(bar, color); no_line(bar); bar.shadow.inherit = False
        add_text(s, x + Inches(0.25), Inches(2.05), w - Inches(0.5), Inches(0.5),
                 title, size=18, bold=True, color=color)
        add_lines(s, x + Inches(0.25), Inches(2.55), w - Inches(0.5), Inches(3.8),
                  [f"•  {it}" for it in items], size=12.5, color=INK_SOFT,
                  line_spacing=1.25, space_after=4)
        x += w + Inches(0.16)

    add_text(s, Inches(0.6), Inches(6.65), Inches(12.28), Inches(0.4),
             "Goal: works on a slow phone, on a small-town network, for everyone.",
             size=12, italic=True, color=MUTED, align=PP_ALIGN.CENTER)


def slide_security(prs):
    s = base_slide(prs, 16, "Security",
                   "Defence in depth, no exceptions, every layer hostile to the layer above")

    items = [
        ("bcrypt cost 12  +  password ≥ 10 chars", "Password hashing tuned for 2026 hardware."),
        ("JWT HS256 in HttpOnly Secure SameSite=Lax cookie", "Server-side sessions table allows revocation."),
        ("CSRF double-submit token (gs_csrf)", "@fastify/csrf-protection on every state-changing POST."),
        ("Parameterised SQL only, zero string concat", "better-sqlite3 prepared statements throughout."),
        ("zod validation at every request boundary", "Body, query, params parsed before any handler runs."),
        ("audit_log on every mutation", "actor, ip, ua, sanitized payload  —  visible in /admin."),
        ("Rate limit 200 req / minute / IP", "via @fastify/rate-limit, /healthz allow-listed."),
        ("Secrets in Google Secret Manager", "Cloud Run mounts JWT_SECRET, RZP_*, DIGILOCKER_*, ADMIN_PASSWORD."),
        ("HSTS, X-Content-Type-Options, Referrer-Policy", "Set globally. HTTPS by default on Cloud Run."),
        ("PII never in logs or URLs", "Stack traces visible only to admin."),
    ]
    x = Inches(0.6)
    y = Inches(1.85)
    rh = Inches(0.50)
    rw = Inches(6.10)
    for i, (k, v) in enumerate(items):
        col = i // 5
        row = i % 5
        cx = x + col * (rw + Inches(0.08))
        cy = y + row * (rh + Inches(0.10))
        card(s, cx, cy, rw, rh, fill=CARD_BG, line_color=LINE, line_w=0.5)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, cy, Inches(0.08), rh)
        set_solid_fill(bar, SAFFRON); no_line(bar); bar.shadow.inherit = False
        add_text(s, cx + Inches(0.20), cy + Inches(0.04), rw - Inches(0.30), Inches(0.25),
                 k, size=12, bold=True, color=INK)
        add_text(s, cx + Inches(0.20), cy + Inches(0.26), rw - Inches(0.30), Inches(0.24),
                 v, size=10.5, color=MUTED)


def slide_observability(prs):
    s = base_slide(prs, 17, "Observability",
                   "Verbose for the admin, silent for the user")

    code_block(
        s, Inches(0.6), Inches(1.90), Inches(12.28), Inches(1.40),
        [
            [{"text": "// pino structured log line"}],
            [{"text": "{ \"level\":30, \"time\":1714660800000, \"reqId\":\"01HX...\","}],
            [{"text": "  \"method\":\"POST\", \"url\":\"/bookings\", \"status\":302, \"ms\":42,"}],
            [{"text": "  \"user_id\":\"01HX...\", \"ip\":\"203.0.113.4\", \"ua\":\"Mozilla/5.0...\" }"}],
        ],
        size=12,
    )

    cards = [
        ("Every request", "ULID req-id, method, url, status, ms, user, ip, ua.", NAVY),
        ("Every mutation", "audit_log row: actor, action, entity, payload (sanitized).", SAFFRON),
        ("/admin SIEM", "Filterable timeline + Server-Sent Events live feed.", GREEN),
        ("Probes", "/healthz returns ok shape. /readyz checks DB + uploads dir.", DEEP),
    ]
    x = Inches(0.6)
    w = Inches(3.0)
    for title, body, color in cards:
        card(s, x, Inches(3.55), w, Inches(2.0), fill=CARD_BG, line_color=color, line_w=1.0)
        add_text(s, x + Inches(0.20), Inches(3.70), w - Inches(0.4), Inches(0.5),
                 title, size=15, bold=True, color=color)
        add_text(s, x + Inches(0.20), Inches(4.20), w - Inches(0.4), Inches(1.7),
                 body, size=12, color=INK_SOFT)
        x += w + Inches(0.08)

    add_text(s, Inches(0.6), Inches(5.80), Inches(12.28), Inches(0.5),
             "User sees friendly error + correlation ID. Super-admin sees stack, "
             "request, payload, and the audit trail — same correlation ID joins them.",
             size=12, italic=True, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(6.30), Inches(12.28), Inches(0.5),
             "Cloud Logging picks pino JSON automatically; alert policy on 5xx ratio.",
             size=12, color=INK_SOFT, align=PP_ALIGN.CENTER)


def slide_deploy(prs):
    s = base_slide(prs, 18, "Cloud Run Deployment",
                   "One container, one shell command, asia-south1 region")

    code_block(
        s, Inches(0.6), Inches(1.85), Inches(6.05), Inches(2.20),
        [
            [{"text": "$ "}, {"text": "./deploy.sh PROJECT_ID", "color": CODE_KEY, "bold": True}],
            [{"text": "  enabling APIs..."}],
            [{"text": "  creating Artifact Registry repo gharsetu-images"}],
            [{"text": "  ensuring 5 secrets in Secret Manager"}],
            [{"text": "  granting roles/secretmanager.secretAccessor"}],
            [{"text": "  submitting cloudbuild.yaml..."}],
            [{"text": "  Service URL: https://gharsetu-...run.app", "color": CODE_STR}],
        ],
        size=12,
    )

    code_block(
        s, Inches(6.85), Inches(1.85), Inches(6.05), Inches(2.20),
        [
            [{"text": "# cloudbuild.yaml steps"}],
            [{"text": "1. docker build  --tag :$SHORT_SHA :latest"}],
            [{"text": "2. push :$SHORT_SHA  to Artifact Registry"}],
            [{"text": "3. push :latest"}],
            [{"text": "4. gcloud run deploy gharsetu"}],
            [{"text": "    --region asia-south1"}],
            [{"text": "    --memory 512Mi  --cpu 1"}],
            [{"text": "    --min-instances 0  --max-instances 10"}],
        ],
        size=12,
    )

    facts = [
        ("~631 MB", "image size"),
        ("~5 s", "cold start"),
        ("₹0", "idle cost"),
        ("80", "concurrency"),
        ("5", "secrets mounted"),
    ]
    x = Inches(0.6)
    w = Inches(2.4)
    for value, label in facts:
        card(s, x, Inches(4.30), w, Inches(1.30), fill=PALE, line_color=SAFFRON, line_w=1.0)
        add_text(s, x, Inches(4.40), w, Inches(0.7),
                 value, size=26, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(5.10), w, Inches(0.5),
                 label, size=12, color=MUTED, align=PP_ALIGN.CENTER)
        x += w + Inches(0.08)

    add_text(s, Inches(0.6), Inches(5.85), Inches(12.28), Inches(0.5),
             "Re-deploy after a code change is the same one command. Idempotent end to end.",
             size=12, italic=True, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(6.30), Inches(12.28), Inches(0.5),
             "JWT_SECRET, RZP_KEY_SECRET, RZP_WEBHOOK_SECRET, DIGILOCKER_CLIENT_SECRET, "
             "ADMIN_PASSWORD  —  all from Secret Manager, never in image.",
             size=11, color=INK_SOFT, align=PP_ALIGN.CENTER, font=CODE_FONT)


def slide_tests(prs):
    s = base_slide(prs, 19, "Testing", "npm test  —  10 smoke checks, all green")

    tests = [
        ("GET /healthz returns ok", True),
        ("GET / serves the home page", True),
        ("GET /search lists Solan listings", True),
        ("GET /api/search?city=Solan returns JSON", True),
        ("Signup + login issues a session cookie", True),
        ("Authenticated student can POST /bookings", True),
        ("POST /pay/order returns an order_id", True),
        ("POST /ondc/v1/search returns ACK", True),
        ("GET /admin without admin cookie is blocked", True),
        ("GET /healthz body shape is { ok: true }", True),
    ]
    x = Inches(0.6)
    w = Inches(6.10)
    rh = Inches(0.46)
    for i, (label, ok) in enumerate(tests):
        col = i // 5
        row = i % 5
        cx = x + col * (w + Inches(0.10))
        cy = Inches(1.85) + row * (rh + Inches(0.10))
        card(s, cx, cy, w, rh, fill=CARD_BG, line_color=LINE, line_w=0.5)
        check = s.shapes.add_shape(MSO_SHAPE.OVAL, cx + Inches(0.12), cy + Inches(0.10), Inches(0.26), Inches(0.26))
        set_solid_fill(check, GREEN); no_line(check); check.shadow.inherit = False
        tf = check.text_frame
        tf.margin_left = Emu(0); tf.margin_right = Emu(0)
        tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = tf.paragraphs[0].add_run()
        r.text = "✓"
        r.font.name = BODY_FONT
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = CARD_BG
        add_text(s, cx + Inches(0.50), cy + Inches(0.10), w - Inches(0.7), Inches(0.30),
                 label, size=12, color=INK)

    # Footer summary
    card(s, Inches(0.6), Inches(5.40), Inches(12.28), Inches(1.35), fill=PALE, line_color=GREEN, line_w=1.0)
    add_text(s, Inches(0.85), Inches(5.50), Inches(11.7), Inches(0.5),
             "10 / 10 passed  —  every prod bug becomes a regression test here.",
             size=15, bold=True, color=GREEN)
    add_text(s, Inches(0.85), Inches(5.85), Inches(11.7), Inches(0.5),
             "Server is spawned per run on an ephemeral port + ephemeral SQLite + ephemeral uploads dir.",
             size=11, color=INK_SOFT, font=CODE_FONT)
    add_text(s, Inches(0.85), Inches(6.20), Inches(11.7), Inches(0.5),
             "npm run lint  =  tsc --noEmit (strict)    npm run build  =  full transpile, runs in <2s.",
             size=11, color=INK_SOFT, font=CODE_FONT)


def slide_next(prs):
    s = base_slide(prs, 20, "What's Next", "After the capstone, before a real launch")

    items = [
        ("Production ONDC subscription",
         "Register with the ONDC registry, generate Ed25519 + X25519 keys, sign every request."),
        ("Cloud SQL for PostgreSQL",
         "Move off ephemeral SQLite. Schema is portable; swap better-sqlite3 for pg."),
        ("Mobile app",
         "PWA first (the SW is already wired); evaluate React Native for native push."),
        ("Owner WhatsApp notifications",
         "Booking and payment events delivered where owners actually read."),
        ("Photo verification",
         "Owner-uploaded shots geo-checked and timestamp-checked before publish."),
        ("Neighbourhood chat",
         "Per-area channels for current renters — keeps the trust loop alive after move-in."),
        ("Payment escrow",
         "Hold the first month + deposit until the renter checks in. Fraud-proofs the deposit."),
    ]
    x = Inches(0.6)
    w = Inches(6.10)
    rh = Inches(0.65)
    for i, (k, v) in enumerate(items):
        col = i // 4
        row = i % 4
        cx = x + col * (w + Inches(0.08))
        cy = Inches(1.85) + row * (rh + Inches(0.20))
        card(s, cx, cy, w, rh, fill=CARD_BG, line_color=LINE, line_w=0.5)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, cy, Inches(0.08), rh)
        set_solid_fill(bar, SAFFRON); no_line(bar); bar.shadow.inherit = False
        add_text(s, cx + Inches(0.22), cy + Inches(0.05), w - Inches(0.30), Inches(0.30),
                 k, size=14, bold=True, color=INK)
        add_text(s, cx + Inches(0.22), cy + Inches(0.32), w - Inches(0.30), Inches(0.30),
                 v, size=11, color=MUTED)


def slide_learnings(prs):
    s = base_slide(prs, 21, "Learnings", "What this capstone taught me")

    cols = [
        ("Engineering",
         ["Strict-mode TypeScript across server + lib.",
          "Native-cloud first: Cloud Run, Secret Manager, Artifact Registry, Cloud Logging.",
          "Protocol design: Beckn / ONDC contract rather than ad-hoc API.",
          "WCAG 2.2 AAA as a build constraint, not an afterthought.",
          "Observability: structured logs + audit trail + SIEM dashboard."],
         NAVY),
        ("Process",
         ["Spec-first: SPEC.md drove every contract and table.",
          "Agent-coordinated parallel workstreams kept progress unblocked.",
          "Ephemeral-first MVP: no infra to manage during the build.",
          "Deploy script idempotent from day one  —  re-run is the workflow.",
          "Smoke suite as the safety net while iterating fast."],
         SAFFRON),
        ("Personal",
         ["Shipping > perfection.  A working slow phone trumps a beautiful demo.",
          "Reading other people's APIs (Razorpay, ONDC, DigiLocker) is half the job.",
          "Real users are not me: i18n + accessibility from day zero.",
          "Document why, not what  —  README is the artefact reviewers actually read.",
          "A capstone is an opening move, not a closing one."],
         GREEN),
    ]
    x = Inches(0.6)
    w = Inches(4.04)
    for title, items, color in cols:
        card(s, x, Inches(1.85), w, Inches(5.0), fill=CARD_BG, line_color=LINE)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.85), w, Inches(0.10))
        set_solid_fill(bar, color); no_line(bar); bar.shadow.inherit = False
        add_text(s, x + Inches(0.25), Inches(2.05), w - Inches(0.5), Inches(0.5),
                 title, size=18, bold=True, color=color)
        add_lines(s, x + Inches(0.25), Inches(2.60), w - Inches(0.5), Inches(4.2),
                  [f"•  {it}" for it in items], size=12, color=INK_SOFT,
                  line_spacing=1.30, space_after=6)
        x += w + Inches(0.16)


def slide_thanks(prs):
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)
    page_background(s)
    accent_bar(s)

    add_text(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(1.4),
             "Thank you.", size=84, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(2.85), Inches(12.1), Inches(0.7),
             "Questions, please.", size=28, italic=True, color=DEEP, align=PP_ALIGN.CENTER)

    div = s.shapes.add_connector(1, Inches(5.0), Inches(3.85), Inches(8.3), Inches(3.85))
    div.line.color.rgb = SAFFRON
    div.line.width = Pt(2.5)

    add_text(s, Inches(0.6), Inches(4.05), Inches(12.1), Inches(0.5),
             "Akshit Thakur",
             size=20, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(4.45), Inches(12.1), Inches(0.5),
             "B.Tech CSE (Cybersecurity)  ·  Shoolini University, Solan",
             size=14, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(4.85), Inches(12.1), Inches(0.5),
             "github.com/akshit-thakur     ·     akshit.thakur@shoolini.edu.in",
             size=14, color=NAVY, align=PP_ALIGN.CENTER, font=CODE_FONT)

    # Tagline at bottom
    card(s, Inches(3.5), Inches(5.85), Inches(6.3), Inches(0.95),
         fill=PALE, line_color=SAFFRON, line_w=1.0)
    add_text(s, Inches(3.5), Inches(5.95), Inches(6.3), Inches(0.5),
             "GharSetu", size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_text(s, Inches(3.5), Inches(6.40), Inches(6.3), Inches(0.4),
             "Apna kamra, apne sheher mein.",
             size=14, italic=True, color=DEEP, align=PP_ALIGN.CENTER, font=HINDI_FONT)

    add_text(s, Inches(0.45), Inches(7.10), Inches(12.4), Inches(0.3),
             "Akshit Thakur  ·  GharSetu Capstone  ·  2026",
             size=10, color=MUTED, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title_page(prs)              # 1
    slide_problem(prs)                 # 2
    slide_insight(prs)                 # 3
    slide_vision(prs)                  # 4
    slide_users(prs)                   # 5
    slide_built(prs)                   # 6
    slide_demo(prs)                    # 7
    slide_architecture(prs)            # 8
    slide_data_model(prs)              # 9
    slide_trust(prs)                   # 10
    slide_ondc(prs)                    # 11
    slide_payments(prs)                # 12
    slide_kyc(prs)                     # 13
    slide_perf(prs)                    # 14
    slide_a11y(prs)                    # 15
    slide_security(prs)                # 16
    slide_observability(prs)           # 17
    slide_deploy(prs)                  # 18
    slide_tests(prs)                   # 19
    slide_next(prs)                    # 20
    slide_learnings(prs)               # 21
    slide_thanks(prs)                  # 22

    prs.save(OUTPUT)
    return prs


def verify(prs):
    size = OUTPUT.stat().st_size
    count = len(prs.slides)
    print(f"file:        {OUTPUT}")
    print(f"size_bytes:  {size}  ({size / 1024:.1f} KB)")
    print(f"slide_count: {count}")
    print("titles:")
    for i, slide in enumerate(prs.slides, 1):
        title_text = ""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                t = "".join(r.text for r in para.runs).strip()
                if t:
                    title_text = t
                    break
            if title_text:
                break
        print(f"  {i:2d}. {title_text}")


if __name__ == "__main__":
    p = build()
    verify(p)
