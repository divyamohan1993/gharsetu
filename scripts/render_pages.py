"""Render the GharSetu capstone .docx and .pptx into HTML fragments.

Re-runnable. Produces three artefacts:

1. src/views/_generated/report-body.html
   - TOC sidebar + <article> body of the report (no <html>/<head> wrapper).
2. src/views/_generated/pitch-body.html
   - 16:9 slide deck markup + control bar.
3. src/public/downloads/<original binaries>
   - Verbatim copies for the on-page download buttons.

Usage:
    python3 scripts/render_pages.py
"""

from __future__ import annotations

import re
import shutil
import sys
import unicodedata
from html import escape as html_escape
from pathlib import Path
from typing import Iterable, List, Optional, Tuple

from docx import Document
from docx.oxml.ns import qn
from docx.text.paragraph import Paragraph
from docx.table import Table

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu


REPO_ROOT = Path(__file__).resolve().parent.parent
REPORT_DOCX = REPO_ROOT / "Akshit_Thakur_Capstone_Report.docx"
PITCH_PPTX = REPO_ROOT / "Akshit_Thakur_Capstone_Presentation.pptx"

VIEWS_GEN = REPO_ROOT / "src" / "views" / "_generated"
DOWNLOADS = REPO_ROOT / "src" / "public" / "downloads"

REPORT_OUT = VIEWS_GEN / "report-body.ejs"
PITCH_OUT = VIEWS_GEN / "pitch-body.ejs"


# ---------------------------------------------------------------------------
# helpers
# ---------------------------------------------------------------------------

_slug_used: dict[str, int] = {}


def slugify(text: str) -> str:
    """Deterministic ASCII slug. Suffix duplicates with -2, -3, ..."""
    norm = unicodedata.normalize("NFKD", text)
    ascii_only = norm.encode("ascii", "ignore").decode("ascii").lower()
    cleaned = re.sub(r"[^a-z0-9]+", "-", ascii_only).strip("-")
    cleaned = re.sub(r"-{2,}", "-", cleaned) or "section"
    n = _slug_used.get(cleaned, 0) + 1
    _slug_used[cleaned] = n
    return cleaned if n == 1 else f"{cleaned}-{n}"


def reset_slug_counter() -> None:
    _slug_used.clear()


def render_runs(paragraph: Paragraph) -> str:
    """Convert paragraph runs to inline HTML, preserving bold/italic."""
    parts: List[str] = []
    for run in paragraph.runs:
        text = run.text
        if not text:
            continue
        encoded = html_escape(text).replace("\n", "<br>")
        if run.bold and run.italic:
            encoded = f"<strong><em>{encoded}</em></strong>"
        elif run.bold:
            encoded = f"<strong>{encoded}</strong>"
        elif run.italic:
            encoded = f"<em>{encoded}</em>"
        parts.append(encoded)
    if not parts and paragraph.text:
        parts.append(html_escape(paragraph.text))
    return "".join(parts)


def iter_block_items(parent) -> Iterable[object]:
    """Yield each top-level Paragraph or Table in document body order."""
    body = parent.element.body
    for child in body.iterchildren():
        if child.tag == qn("w:p"):
            yield Paragraph(child, parent)
        elif child.tag == qn("w:tbl"):
            yield Table(child, parent)


# ---------------------------------------------------------------------------
# report
# ---------------------------------------------------------------------------


def render_report() -> Tuple[str, int]:
    reset_slug_counter()
    doc = Document(str(REPORT_DOCX))

    toc_entries: List[Tuple[str, str]] = []  # (slug, title)
    body_parts: List[str] = []

    for block in iter_block_items(doc):
        if isinstance(block, Paragraph):
            body_parts.append(render_paragraph(block, toc_entries))
        elif isinstance(block, Table):
            body_parts.append(render_table(block))

    article = (
        '<article class="report__body" lang="en">\n'
        + "\n".join(p for p in body_parts if p)
        + "\n</article>"
    )

    toc_items = "\n".join(
        f'      <li><a href="#{slug}">{html_escape(title)}</a></li>'
        for slug, title in toc_entries
    )
    toc = (
        '<nav class="report__toc" aria-label="Table of contents">\n'
        '  <p class="report__toc-label">Contents</p>\n'
        f'  <ol class="report__toc-list">\n{toc_items}\n  </ol>\n'
        "</nav>"
    )

    return toc + "\n" + article, len(toc_entries)


def render_paragraph(p: Paragraph, toc_entries: List[Tuple[str, str]]) -> str:
    style = (p.style.name or "").strip()
    text = p.text.strip()
    inner = render_runs(p)

    if not text and not inner:
        return ""

    if style == "Heading 1":
        slug = slugify(text)
        toc_entries.append((slug, text))
        return f'<h2 class="report__h1" id="{slug}">{inner or html_escape(text)}</h2>'
    if style == "Heading 2":
        return f'<h3 class="report__h2">{inner or html_escape(text)}</h3>'
    if style == "Heading 3":
        return f'<h4 class="report__h3">{inner or html_escape(text)}</h4>'
    if style.startswith("List Bullet"):
        return f'<ul class="report__ul"><li>{inner}</li></ul>'
    if style.startswith("List Number"):
        return f'<ol class="report__ol"><li>{inner}</li></ol>'

    return f'<p class="report__p">{inner}</p>'


def render_table(table: Table) -> str:
    rows = list(table.rows)
    if not rows:
        return ""
    head_cells = "".join(
        f'<th scope="col">{html_escape(cell.text.strip())}</th>'
        for cell in rows[0].cells
    )
    body_rows: List[str] = []
    for row in rows[1:]:
        cells = "".join(
            f"<td>{html_escape(cell.text.strip()).replace(chr(10), '<br>')}</td>"
            for cell in row.cells
        )
        body_rows.append(f"<tr>{cells}</tr>")
    body_html = "\n".join(body_rows)
    return (
        '<div class="report__table-wrap">'
        '<table class="report__table">'
        f"<thead><tr>{head_cells}</tr></thead>"
        f"<tbody>\n{body_html}\n</tbody>"
        "</table>"
        "</div>"
    )


# ---------------------------------------------------------------------------
# pitch
# ---------------------------------------------------------------------------


def emu_to_pct(value: int, total: int) -> float:
    if total <= 0:
        return 0.0
    return round((value / total) * 10000) / 100


def pt_to_px(pt_value: float) -> float:
    return round(pt_value * 1.333 * 100) / 100


def safe_color_hex(run) -> Optional[str]:
    try:
        rgb = run.font.color.rgb
    except Exception:
        return None
    if rgb is None:
        return None
    return f"#{str(rgb)}"


def render_text_frame(tf, default_align: Optional[str] = None) -> str:
    parts: List[str] = []
    for para in tf.paragraphs:
        align = default_align
        if para.alignment is not None:
            align_map = {1: "left", 2: "center", 3: "right", 4: "justify"}
            align = align_map.get(int(para.alignment), align)
        run_html_parts: List[str] = []
        for run in para.runs:
            text = run.text
            if not text:
                continue
            encoded = html_escape(text)
            styles: List[str] = []
            if run.font.size is not None:
                size_pt = run.font.size.pt
                styles.append(f"font-size:{pt_to_px(size_pt)}px")
            color = safe_color_hex(run)
            if color is not None:
                styles.append(f"color:{color}")
            tag = encoded
            if run.font.bold and run.font.italic:
                tag = f"<strong><em>{tag}</em></strong>"
            elif run.font.bold:
                tag = f"<strong>{tag}</strong>"
            elif run.font.italic:
                tag = f"<em>{tag}</em>"
            if styles:
                tag = f'<span style="{";".join(styles)}">{tag}</span>'
            run_html_parts.append(tag)
        if not run_html_parts and para.text:
            run_html_parts.append(html_escape(para.text))
        attrs = ""
        if align:
            attrs = f' style="text-align:{align}"'
        parts.append(f"<p{attrs}>{''.join(run_html_parts) or '&nbsp;'}</p>")
    return "\n".join(parts)


def shape_fill_hex(shape) -> Optional[str]:
    try:
        fill = shape.fill
        if fill.type != 1:  # SOLID
            return None
        rgb = fill.fore_color.rgb
        if rgb is None:
            return None
        return f"#{str(rgb)}"
    except Exception:
        return None


def is_title_shape(shape, idx: int, title_seen: bool) -> bool:
    try:
        ph_type = shape.placeholder_format.type
        if ph_type is not None and "TITLE" in str(ph_type):
            return True
    except Exception:
        pass
    if title_seen:
        return False
    if not shape.has_text_frame:
        return False
    text = shape.text_frame.text.strip()
    if not text or len(text) > 120:
        return False
    return idx <= 4


def render_pitch() -> Tuple[str, int]:
    pres = Presentation(str(PITCH_PPTX))
    slide_w = pres.slide_width
    slide_h = pres.slide_height

    slides_html: List[str] = []
    for i, slide in enumerate(pres.slides):
        slide_num = i + 1
        slide_id = f"slide-{slide_num}"
        slide_title = ""

        body_parts: List[str] = []
        title_seen = False

        for idx, shape in enumerate(slide.shapes):
            left = max(0, shape.left or 0)
            top = max(0, shape.top or 0)
            width = max(0, shape.width or 0)
            height = max(0, shape.height or 0)
            l_pct = emu_to_pct(left, slide_w)
            t_pct = emu_to_pct(top, slide_h)
            w_pct = emu_to_pct(width, slide_w)
            h_pct = emu_to_pct(height, slide_h)

            if shape.has_text_frame and shape.text_frame.text.strip():
                txt = shape.text_frame.text.strip()
                if not title_seen and is_title_shape(shape, idx, title_seen):
                    title_seen = True
                    slide_title = txt
                    body_parts.append(
                        f'<h2 class="slide__title" id="{slide_id}-title" '
                        f'style="position:absolute;left:{l_pct}%;top:{t_pct}%;'
                        f'width:{w_pct}%;height:{h_pct}%;">{html_escape(txt)}</h2>'
                    )
                    continue
                fill_bg = shape_fill_hex(shape)
                style_extras: List[str] = []
                if fill_bg is not None:
                    style_extras.append(f"background:{fill_bg}")
                    style_extras.append("padding:0.4em 0.6em")
                    style_extras.append("border-radius:6px")
                style_attr = (
                    f"position:absolute;left:{l_pct}%;top:{t_pct}%;"
                    f"width:{w_pct}%;height:{h_pct}%;"
                ) + (";".join(style_extras) + ";" if style_extras else "")
                inner = render_text_frame(shape.text_frame)
                body_parts.append(
                    f'<div class="slide__text" style="{style_attr}">{inner}</div>'
                )
            else:
                fill_bg = shape_fill_hex(shape)
                if fill_bg is None:
                    continue
                if width <= 0 or height <= 0:
                    continue
                if l_pct == 0 and t_pct == 0 and w_pct >= 99 and h_pct >= 99:
                    body_parts.append(
                        f'<div class="slide__shape slide__bg" '
                        f'style="position:absolute;inset:0;background:{fill_bg};"></div>'
                    )
                    continue
                body_parts.append(
                    f'<div class="slide__shape" '
                    f'style="position:absolute;left:{l_pct}%;top:{t_pct}%;'
                    f'width:{w_pct}%;height:{h_pct}%;background:{fill_bg};'
                    f'border-radius:6px;"></div>'
                )

        if not slide_title:
            slide_title = f"Slide {slide_num}"
            body_parts.insert(
                0,
                f'<h2 class="slide__title visually-hidden" id="{slide_id}-title">{html_escape(slide_title)}</h2>',
            )

        active_class = " is-active" if slide_num == 1 else ""
        section = (
            f'<section class="slide{active_class}" id="{slide_id}" '
            f'data-slide="{slide_num}" '
            f'data-slide-title="{html_escape(slide_title)}" '
            f'aria-labelledby="{slide_id}-title">\n'
            + "\n".join(body_parts)
            + "\n</section>"
        )
        slides_html.append(section)

    n = len(slides_html)
    deck = (
        f'<div class="pitch" id="pitch" data-slide-count="{n}" tabindex="-1" '
        f'aria-roledescription="carousel" aria-label="Capstone slide deck">\n'
        '  <div class="pitch__viewport">\n'
        + "\n".join(slides_html)
        + "\n  </div>\n</div>\n"
        '<noscript><style>.pitch .slide{display:block !important;position:relative !important;'
        'inset:auto !important;margin-bottom:2rem;border:1px solid var(--c-border-soft);}</style></noscript>\n'
        '<div class="pitch__controls" aria-label="Slide navigation">\n'
        '  <button type="button" class="btn btn--ghost" data-slide-action="prev" aria-label="Previous slide">&larr; Prev</button>\n'
        f'  <span class="pitch__counter" aria-live="polite" aria-atomic="true">1 / {n}</span>\n'
        '  <button type="button" class="btn btn--ghost" data-slide-action="next" aria-label="Next slide">Next &rarr;</button>\n'
        '  <a class="btn btn--primary" href="/pitch.pptx" download>Download .pptx</a>\n'
        '</div>\n'
        '<p class="pitch__help"><kbd>&larr;</kbd> <kbd>&rarr;</kbd> navigate &middot; '
        '<kbd>F</kbd> fullscreen &middot; <kbd>Esc</kbd> exit</p>\n'
    )
    return deck, n


# ---------------------------------------------------------------------------
# main
# ---------------------------------------------------------------------------


def main() -> int:
    if not REPORT_DOCX.exists():
        print(f"[render_pages] missing report: {REPORT_DOCX}", file=sys.stderr)
        return 1
    if not PITCH_PPTX.exists():
        print(f"[render_pages] missing pitch: {PITCH_PPTX}", file=sys.stderr)
        return 1

    VIEWS_GEN.mkdir(parents=True, exist_ok=True)
    DOWNLOADS.mkdir(parents=True, exist_ok=True)

    report_html, h_count = render_report()
    REPORT_OUT.write_text(report_html, encoding="utf-8")

    pitch_html, slide_count = render_pitch()
    PITCH_OUT.write_text(pitch_html, encoding="utf-8")

    report_dl = DOWNLOADS / REPORT_DOCX.name
    pitch_dl = DOWNLOADS / PITCH_PPTX.name
    shutil.copy2(REPORT_DOCX, report_dl)
    shutil.copy2(PITCH_PPTX, pitch_dl)

    print("[render_pages] artefacts:")
    for path in (REPORT_OUT, PITCH_OUT, report_dl, pitch_dl):
        size = path.stat().st_size
        print(f"  {path.relative_to(REPO_ROOT)}  ({size:,} bytes)")
    print(f"[render_pages] headings (H1) in report: {h_count}")
    print(f"[render_pages] slides in pitch: {slide_count}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
